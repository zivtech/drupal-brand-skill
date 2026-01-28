#!/usr/bin/env python3
"""
Drupal Brand Compliance Analyzer for PowerPoint Presentations

Scans a PPTX file and identifies slides that need manual intervention
for brand compliance (hard-coded fonts, off-brand colors).

Usage:
    python analyze_deck.py path/to/presentation.pptx

Requirements:
    pip install python-pptx
"""

import sys
import os

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
except ImportError:
    print("Error: python-pptx is required. Install with: pip install python-pptx")
    sys.exit(1)

# Drupal brand colors (lowercase hex for comparison)
DRUPAL_COLORS = {
    '009cde': 'Drupal Blue',
    '006aa9': 'Drupal Dark Blue',
    '12285f': 'Drupal Navy',
    'ccedf9': 'Drupal Light Blue',
    'ffc423': 'Drupal Yellow',
    'f46351': 'Drupal Red',
    '397618': 'Drupal Green',
    'ccbaf4': 'Drupal Purple',
    '000000': 'Black',
    'ffffff': 'White',
}

# Fonts that should be replaced
BAD_FONTS = ['ubuntu', 'arial', 'calibri', 'helvetica', 'verdana', 'tahoma', 'times', 'georgia']
GOOD_FONTS = ['zt gatha', 'noto sans', 'noto']


def rgb_to_hex(rgb):
    """Convert RGBColor to hex string."""
    if rgb is None:
        return None
    return f'{rgb.red:02x}{rgb.green:02x}{rgb.blue:02x}'.lower()


def is_brand_color(hex_color):
    """Check if color is a Drupal brand color."""
    if hex_color is None:
        return True  # Theme colors are OK
    return hex_color.lower() in DRUPAL_COLORS


def is_bad_font(font_name):
    """Check if font needs replacement."""
    if font_name is None:
        return False
    font_lower = font_name.lower()
    # Skip if it's already a good font
    if any(good in font_lower for good in GOOD_FONTS):
        return False
    # Flag if it's a known bad font
    return any(bad in font_lower for bad in BAD_FONTS)


def analyze_shape(shape):
    """Analyze a shape for brand compliance issues."""
    issues = []

    # Check text frames
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                # Check font
                if run.font.name and is_bad_font(run.font.name):
                    issues.append(f"Hard-coded font: {run.font.name}")

                # Check font color
                try:
                    if run.font.color.type is not None and run.font.color.rgb:
                        hex_color = rgb_to_hex(run.font.color.rgb)
                        if not is_brand_color(hex_color):
                            issues.append(f"Off-brand text color: #{hex_color}")
                except:
                    pass

    # Check fill colors
    if hasattr(shape, 'fill'):
        try:
            if shape.fill.type is not None:
                if hasattr(shape.fill, 'fore_color') and shape.fill.fore_color.rgb:
                    hex_color = rgb_to_hex(shape.fill.fore_color.rgb)
                    if not is_brand_color(hex_color):
                        issues.append(f"Off-brand fill color: #{hex_color}")
        except:
            pass

    # Check line colors
    if hasattr(shape, 'line'):
        try:
            if shape.line.color.type is not None and shape.line.color.rgb:
                hex_color = rgb_to_hex(shape.line.color.rgb)
                if not is_brand_color(hex_color):
                    issues.append(f"Off-brand line color: #{hex_color}")
        except:
            pass

    return list(set(issues))  # Remove duplicates


def get_slide_title(slide):
    """Extract slide title if present."""
    if slide.shapes.title:
        title = slide.shapes.title.text.strip()
        return title[:50] + "..." if len(title) > 50 else title
    # Try to find first text
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            text = shape.text_frame.text.strip()
            return text[:50] + "..." if len(text) > 50 else text
    return "(No title)"


def analyze_presentation(pptx_path):
    """Analyze entire presentation for brand compliance."""
    print(f"\nAnalyzing: {pptx_path}")
    print("=" * 70)

    prs = Presentation(pptx_path)

    slides_with_issues = []
    all_fonts = set()

    for i, slide in enumerate(prs.slides, 1):
        slide_issues = []

        for shape in slide.shapes:
            issues = analyze_shape(shape)
            slide_issues.extend(issues)

            # Collect fonts for summary
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            all_fonts.add(run.font.name)

        if slide_issues:
            slides_with_issues.append({
                'num': i,
                'title': get_slide_title(slide),
                'issues': list(set(slide_issues))
            })

    # Categorize issues
    font_issues = []
    color_issues = []
    both_issues = []

    for slide in slides_with_issues:
        has_font = any('font' in i.lower() for i in slide['issues'])
        has_color = any('color' in i.lower() for i in slide['issues'])

        if has_font and has_color:
            both_issues.append(slide)
        elif has_font:
            font_issues.append(slide)
        elif has_color:
            color_issues.append(slide)

    # Print results
    print(f"\nTotal slides: {len(prs.slides)}")
    print(f"Clean slides: {len(prs.slides) - len(slides_with_issues)}")
    print(f"Slides needing fixes: {len(slides_with_issues)}")

    if all_fonts:
        bad_fonts_found = [f for f in all_fonts if is_bad_font(f)]
        if bad_fonts_found:
            print(f"\nNon-brand fonts found: {', '.join(sorted(bad_fonts_found))}")

    if both_issues:
        print(f"\n{'='*70}")
        print(f"HIGH PRIORITY - Font AND Color Issues ({len(both_issues)} slides)")
        print("-" * 50)
        for slide in both_issues:
            print(f"\nSlide {slide['num']}: {slide['title']}")
            for issue in sorted(set(slide['issues'])):
                print(f"  - {issue}")

    if font_issues:
        print(f"\n{'='*70}")
        print(f"MEDIUM PRIORITY - Font Issues ({len(font_issues)} slides)")
        print("-" * 50)
        for slide in font_issues:
            print(f"\nSlide {slide['num']}: {slide['title']}")
            for issue in sorted(set(slide['issues'])):
                print(f"  - {issue}")

    if color_issues:
        print(f"\n{'='*70}")
        print(f"LOWER PRIORITY - Color Issues ({len(color_issues)} slides)")
        print("-" * 50)
        for slide in color_issues:
            print(f"\nSlide {slide['num']}: {slide['title']}")
            for issue in sorted(set(slide['issues'])):
                print(f"  - {issue}")

    # Summary
    print(f"\n{'='*70}")
    print("SUMMARY")
    print("=" * 70)
    print(f"Total slides: {len(prs.slides)}")
    print(f"Clean (no issues): {len(prs.slides) - len(slides_with_issues)}")
    print(f"Need attention: {len(slides_with_issues)}")
    print(f"  - High priority (font+color): {len(both_issues)}")
    print(f"  - Medium priority (font): {len(font_issues)}")
    print(f"  - Lower priority (color): {len(color_issues)}")

    if slides_with_issues:
        print(f"\nSlides needing fixes: {', '.join(str(s['num']) for s in slides_with_issues)}")

    # Recommended actions
    print(f"\n{'='*70}")
    print("RECOMMENDED ACTIONS")
    print("=" * 70)

    bad_fonts_found = [f for f in all_fonts if is_bad_font(f)]
    if bad_fonts_found:
        print("\n1. BULK FONT REPLACEMENT (PowerPoint: Home → Replace → Replace Fonts)")
        for font in sorted(bad_fonts_found):
            if 'light' in font.lower():
                print(f"   {font} → ZT Gatha")
            else:
                print(f"   {font} → Noto Sans")

    if color_issues or both_issues:
        print("\n2. FIX OFF-BRAND COLORS")
        print("   Review flagged slides and update to Drupal palette:")
        print("   Blue: #009CDE | Navy: #12285F | Yellow: #FFC423")

    print("\n3. RUN THIS SCRIPT AGAIN to verify fixes")

    return len(slides_with_issues)


def main():
    if len(sys.argv) < 2:
        print("Usage: python analyze_deck.py <presentation.pptx>")
        print("\nAnalyzes a PowerPoint file for Drupal brand compliance issues.")
        sys.exit(1)

    pptx_path = sys.argv[1]

    if not os.path.exists(pptx_path):
        print(f"Error: File not found: {pptx_path}")
        sys.exit(1)

    if not pptx_path.lower().endswith('.pptx'):
        print("Error: File must be a .pptx PowerPoint file")
        sys.exit(1)

    issues_count = analyze_presentation(pptx_path)
    sys.exit(0 if issues_count == 0 else 1)


if __name__ == "__main__":
    main()
